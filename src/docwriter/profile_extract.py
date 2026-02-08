from __future__ import annotations

import io
from typing import List

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore

from .rfp_extract import extract_docx_text, extract_pdf_text


def extract_pptx_text(data_bytes: bytes) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is required to extract PPTX text")
    deck = Presentation(io.BytesIO(data_bytes))
    parts: List[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def extract_profile_text(data_bytes: bytes, extension: str) -> str:
    ext = extension.lower().lstrip(".")
    if ext == "pdf":
        return extract_pdf_text(data_bytes)
    if ext == "docx":
        return extract_docx_text(data_bytes)
    if ext == "pptx":
        return extract_pptx_text(data_bytes)
    raise ValueError(f"Unsupported profile file extension: {ext}")
